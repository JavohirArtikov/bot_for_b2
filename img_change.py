from PIL import Image

img = Image.open("bot_for_b\images\searchm.jpg")

new_size = (672, 513)

resized_img = img.resize(new_size)

resized_img.save("bot_for_b\images\searchm.jpg")


# import pytesseract
# from PIL import Image

# # Открыть изображение
# img = Image.open("bot_for_b\images\welcome.jpg")

# # Преобразовать изображение в черно-белый формат
# img = img.convert('L')

# # Извлечь текст из изображения
# text = pytesseract.image_to_string(img)

# # Вывести результат
# print(text)


# import easyocr

# # Создать объект Reader
# def easyocr_recognition(path_img):
#     return easyocr.Reader(["ru"]).readtext(path_img, detail=0, paragraph=True, text_threshold=0.8)


# print(easyocr_recognition("bot_for_b\\images\\translate.jpg"))

# image = Image.open("D:\\Kwork_words\\bot_for_b\\images\\translate.jpg")


# from PIL import Image
# import pytesseract

# pytesseract.pytesseract.tesseract_cmd = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"

# image = Image.open("D:\mywords09.01.2022\Python Задачки\photo_2024-02-04_11-40-57.jpg")

# text = pytesseract.image_to_string(image, lang="eng")

# # print(text)
# from selenium import webdriver
# from selenium.webdriver.chrome.options import Options
# from selenium.webdriver.common.by import By
# from selenium.webdriver.chrome.service import Service
# from selenium.webdriver.support.ui import WebDriverWait
# from selenium.webdriver.support import expected_conditions as EC
# from selenium.webdriver.common.action_chains import ActionChains
# from selenium.webdriver.common.keys import Keys
# from bs4 import BeautifulSoup
# import time
# from urllib.parse import urljoin


# def search_media_files(msg):
#     options = Options()
#     # options.add_argument("--headless")
#     options.add_argument("--lang=ru-RU")
#     options.add_experimental_option("excludeSwitches", ["enable-logging"])
#     options.add_argument(f"userAgent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36")

#     driver = webdriver.Chrome(options=options)
#     driver.get("https://kinogo.biz/")

#     time.sleep(2)
#     search_input = driver.find_element(By.XPATH, "//input[@placeholder='Поиск...']")
#     search_input.send_keys(msg)

#     time.sleep(1)
#     search_input.send_keys(Keys.ENTER)

#     time.sleep(3)

#     soup = BeautifulSoup(driver.page_source, "html.parser")
#     names = soup.find_all("div", class_="shortstory__title")

#     description_text = "🔍В нашей базе нашлось несколько <b>Фильмов</b>, <b>Сериалов</b>, <b>Мультфильмов</b> с подобным названием:\n\n"
#     if names:
#         for number, name in zip(range(1, len(names)+1), names):
#             description_text += f"🔺[{number}]{name.text}\n"

#         description_text += "\n✍️ Укажите цифру <b>Фильма</b>, <b>Сериала</b>, <b>Мультфильма</b>, который вам нужен!"
#         return description_text

#     else:
#         description_text = "Ничего не нашлось!"
#         return description_text


# print(search_media_files("мстители"))
# Путь к изображению с текстом


# import vk_api
# import os

# Введите ваши данные VK
# ac_token="vk1.a.kp_loUO8fn74pwrcZmNa8ynnrjj5utJdl-fu07wKFuMpWgqK9Ltw3Nd948cnrPpzFnL_Gq5zUX2t3-2F-hbEcSUNZo_Yt7zxK-Cqd1uyBG_XVevHMRScB6ecc_Y_CEiRT6LFba75Kf_I1GmQI369_9X-868NKzG1G4GKbcajkEfITiRWAIlokuqWau-NtWpi50dAugGQFfcSsRC3vJziEA"

# import vk_api
# print(vk_api.__version__)
# # Ваш токен для доступа к API VK
# token = 'vk1.a.kp_loUO8fn74pwrcZmNa8ynnrjj5utJdl-fu07wKFuMpWgqK9Ltw3Nd948cnrPpzFnL_Gq5zUX2t3-2F-hbEcSUNZo_Yt7zxK-Cqd1uyBG_XVevHMRScB6ecc_Y_CEiRT6LFba75Kf_I1GmQI369_9X-868NKzG1G4GKbcajkEfITiRWAIlokuqWau-NtWpi50dAugGQFfcSsRC3vJziEA'

# # Авторизация
# vk_session = vk_api.VkApi(token=token)
# vk = vk_session.get_api()


# a = vk_api.audio.VkAudio(vk)
# b = a.search(q="rap god")
# print(b)
# Функция для поиска песни по названию и получения её ID
# def find_song_id_by_title(title):
#     a
#     response = vk.audio.search(q=title, count=1)
#     if response['count'] > 0:
#         song = response['items'][0]
#         song_id = song['id']
#         return song_id
#     else:
#         return None

# # Пример использования
# song_title = "Rap God"
# song_id = find_song_id_by_title(song_title)
# if song_id:
#     print(f"Идентификатор песни с названием '{song_title}': {song_id}")
# else:
#     print(f"Песня с названием '{song_title}' не найдена.")



# import spotipy
# from spotipy.oauth2 import SpotifyClientCredentials
# import os

# # Настройки клиента Spotify
# client_id = '11d5f379729943328d47e96c278700f2'
# client_secret = '62bebe7596f14192851c50570c2e570b'

# sp = spotipy.Spotify(client_credentials_manager=SpotifyClientCredentials(client_id=client_id, client_secret=client_secret))

# # Функция для поиска трека и получения ID
# def find_track_id(track_name):
#     results = sp.search(q=track_name, type='track')
    
#     if results['tracks']['items']:
#         track_id = results['tracks']['items'][0]['id']
#         return track_id
#     else:
#         return None

# # Пример использования
# track_name = "Till I Colldapse"
# track_id = find_track_id(track_name)

# if track_id:
#     print(f"ID трека '{track_name}': {track_id}")
# else:
#     print(f"Трек '{track_name}' не найден.")



    # print(audio_file)
    # await message.answer_audio(open(audio_file, 'rb'))



# from pydeezer import Deezer

# dz = Deezer()

# song_name = "No Lie"
# results = dz.search_tracks(song_name)

# if results:
#     track = results[0]
#     track_name = track['title']
#     artist = track['artist']['name']
#     preview_url = track['preview']

#     print(preview_url)

# else:
#     print(results)


# import soundcloud

# client = soundcloud.Client(client_id="2876adf40e30f73eb3cfe0af77f3a77d", client_secret="80fc6747b2a92ea81b9d29455a5df35a")

# song_name = "Rap God"
# results = client.get("/tracks", q=song_name)

# if results:
#     track = results[0]
#     track_url = track.permalink_url

#     print(track_url)
# else:
#     print("Песня не найдена.")

# from googleapiclient.discovery import build
# from pytube import YouTube
# import youtube_dl
# import os
# # import json

# # Замените "YOUR_API_KEY" на ваш собственный ключ API
# api_key = "AIzaSyCmbgsicv8pOyPTUzo7DsmQASoIcIBpUZU"

# # Создаем объект YouTube API
# youtube = build('youtube', 'v3', developerKey=api_key)

# # Функция для поиска видео по названию песни
# def search_song_video(song_name):
#     search_response = youtube.search().list(
#         q=song_name,
#         part='id,snippet',
#         maxResults=1
#     ).execute()

#     video_id = search_response['items'][0]['id']['videoId']
    
#     return "https://www.youtube.com/watch?v=" + video_id
    

# song_name = input("Введите название песни: ")
# video_id = search_song_video(song_name)


# print(video_id)

# from vk_api import VkApi
# from vk_api.audio import VkAudio

# login = "+998914314489"
# password = "914314489"

# login="996505787504"
# password="rxh12h17Y"

# vk_session = VkApi( 
#   login=login,
#   password=password,
# #   token="vk1.a.kp_loUO8fn74pwrcZmNa8ynnrjj5utJdl-fu07wKFuMpWgqK9Ltw3Nd948cnrPpzFnL_Gq5zUX2t3-2F-hbEcSUNZo_Yt7zxK-Cqd1uyBG_XVevHMRScB6ecc_Y_CEiRT6LFba75Kf_I1GmQI369_9X-868NKzG1G4GKbcajkEfITiRWAIlokuqWau-NtWpi50dAugGQFfcSsRC3vJziEA",
#   api_version='5.81',
#   app_id=2685278
# )

# vkU = vk_session.auth()

# vk_audio = VkAudio(vkU)
# q = "Rap God"
# audio = next(vk_audio.search_iter(q=q))
# url = audio['url']

# print(url)

# import vkpymusic

# acsess = vkaudiotoken.


# import requests

# # Задаем параметры запроса
# params = {
#     "q": "Imagine",
#     "limit": 10,
# }

# # Отправляем запрос
# response = requests.get("https://mp3iq.net/api/search", params=params)

# # Получаем результаты
# results = response.json()

# # Печатаем названия песен
# for song in results["songs"]:
#     print(song["title"])



# from selenium import webdriver
# from selenium.webdriver.chrome.options import Options
# from selenium.webdriver.common.by import By
# from selenium.webdriver.common.keys import Keys
# from bs4 import BeautifulSoup
# import time
# from urllib.parse import urljoin


# def get_music(title_music):
#     options = Options()
#     options.add_argument("--headless")
#     options.add_argument("--lang=ru-RU")
#     options.add_experimental_option("excludeSwitches", ["enable-logging"])
#     options.add_argument("userAgent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36")

#     driver = webdriver.Chrome(options=options)
#     driver.get("https://megapesni.com/")
#     time.sleep(1)

#     search_input = driver.find_element(By.CLASS_NAME, "form-search.ui-autocomplete-input")
#     search_input.send_keys(title_music)
#     search_input.send_keys(Keys.ENTER)
#     time.sleep(2)

#     soup = BeautifulSoup(driver.page_source, "html.parser")
#     link_for_audio = soup.find_all("a", class_="popular-play__item")
#     names = soup.find_all("div", class_="popular-play-name")

#     print(len(link_for_audio))
#     for link, name, num_songs in zip(link_for_audio, names, range(1, 11)):
#         music_link = link["data-url"]
#         music_name = name.text
#         print(f"Ссылка песни: {music_link}\nИмя песни: {music_name}\n\n")


# get_music("Lose yourself")

# import mysql.connector

# connection = mysql.connector.connect(
#     host="localhost",
#     user="root",
#     password="914314489",
#     database="mult_function_bot"
# )

# cursor = connection.cursor()
# # cursor.execute("UPDATE tgbot_multifunction SET first_song = 'adfsfd' where first_song = 0")
# # cursor.execute("UPDATE tgbot_multifunction SET first_song = %s, second_song = %s, third_song = %s, fourth_song = %s, fiveth_song = %s, sixth_song = %s, seventh_song = %s, eighth_song = %s, nineth_song = %s, tenth_song = %s WHERE id = 1;", ('f', 'f2', 'f3', 'f4', 'f5', 'f6', 'f7', 'f8', 'f9', 'f10'))
# cursor.execute("DELETE FROM tgbot_multifunction;")
# # cursor.execute("INSERT INTO tgbot_multifunction (user_id) VALUES (%s);", ([12345]))
# connection.commit()


# import cv2
# import numpy as np

# def enhance_image(image_path):
#     # Загрузка изображения
#     image = cv2.imread(image_path)
    
#     # Коррекция перспективы
#     gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
#     corners = cv2.goodFeaturesToTrack(gray, 4, 0.01, 10)
#     corners = np.float32(corners)
#     new_corners = np.array([[0, 0], [0, image.shape[0]], [image.shape[1], 0], [image.shape[1], image.shape[0]]], dtype=np.float32)
#     M = cv2.getPerspectiveTransform(corners, new_corners)
#     image = cv2.warpPerspective(image, M, (image.shape[1], image.shape[0]))

#     # Улучшение контраста и яркости
#     image = cv2.convertScaleAbs(image, alpha=1.5, beta=50)
    
#     # Удаление шумов
#     image = cv2.fastNlMeansDenoisingColored(image, None, 10, 10, 7, 21)
    
#     return image

# # Путь к изображению
# image_path = "D:\mywords09.01.2022\Python Задачки\photo_2024-02-04_11-38-51.jpg"

# # Улучшение изображения
# enhanced_image = enhance_image(image_path)

# # Сохранение улучшенного изображения
# cv2.imwrite("D:\mywords09.01.2022\Python Задачки\enhanced_image.jpg", enhanced_image)


# from scihub.
# https://ru.pdfdrive.com/



# from selenium import webdriver
# from selenium.webdriver.chrome.options import Options
# from selenium.webdriver.common.by import By
# from selenium.webdriver.common.keys import Keys
# from bs4 import BeautifulSoup
# import time
# from urllib.parse import urljoin


# def get_book(title_book):
#     options = Options()
#     # options.add_argument("--headless")
#     options.add_argument("--lang=ru-RU")
#     options.add_experimental_option("excludeSwitches", ["enable-logging"])
#     options.add_argument("userAgent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36")

#     driver = webdriver.Chrome(options=options)
#     driver.get("https://avidreaders.ru/")
#     time.sleep(1)

#     search_input = driver.find_element(By.CSS_SELECTOR, "input[id='search-txt']")
#     search_input.send_keys(title_book)
#     search_input.send_keys(Keys.ENTER)
#     time.sleep(1)

#     soup = BeautifulSoup(driver.page_source, "html.parser")
#     link_for_book = soup.find_all("a", class_="btn")
#     names = soup.find_all("div", class_="book_name")

#     print(len(link_for_book))
#     for link, name, num_songs in zip(link_for_book, names, range(1, 11)):
#         book_link = link["href"]
#         book_name = name.text
#         print(f"Ссылка книги: {book_link}\nИмя книги: {book_name}\n\n")


# get_book("Марк  дворецкий")


# import requests

# def download_insta_func(url_video):

#     url = "https://instagram-post-and-reels-downloader.p.rapidapi.com/insta/"

#     querystring = {"url":url_video}

#     headers = {
#         "X-RapidAPI-Key": "51efbfb617mshfe027eb5782a110p14e4a0jsn11c17cc1bd33",
#         "X-RapidAPI-Host": "instagram-post-and-reels-downloader.p.rapidapi.com"
#     }

#     response = requests.get(url, headers=headers, params=querystring)

#     data = response.json()
#     # insta_video = data['detail']['data']['items']['0']['urls']['0']['urlWrapped']
#     # insta_video = data['detail']['data']['items']['0']['urls']['0']['urlWrapped']

#     return data



# print(download_insta_func("https://www.instagram.com/reels/C07Faj1OpPb/"))


# import webbrowser

# webbrowser.open("https://www.instagram.com/reels/C2qYNnGPcox/")


# import requests

# def get_youtube_video(id_y_video):
#     url = "https://ytstream-download-youtube-videos.p.rapidapi.com/dl"

#     querystring = {"id":id_y_video}

#     headers = {
#         "X-RapidAPI-Key": "51efbfb617mshfe027eb5782a110p14e4a0jsn11c17cc1bd33",
#         "X-RapidAPI-Host": "ytstream-download-youtube-videos.p.rapidapi.com"
#     }

#     response = requests.get(url, headers=headers, params=querystring)

#     y_video = response.json()
#     video_url = y_video["formats"][0]['url']

#     return video_url

# print(get_youtube_video("39GvuDAacqs"))

# import re

# url = "https://www.youtube.com/watch?v=kTsxLyvL_NY"

# # Извлечение текста после "="
# match = re.search(r"=([^&]*)", url)

# # Текст после "="
# text = match.group(1)

# # Вывод текста
# print(text)


# import re

# # URL видео
# url = "https://www.youtube.com/shorts/kTsxLyvL_NY"

# # Извлечение текста после "https://www.youtube.com/shorts/"
# match = re.search(r"https://www.youtube.com/shorts/(.*)", url)

# # Текст после "https://www.youtube.com/shorts/"
# text = match.group(1)

# # Вывод текста
# print(text)


import openai

OPENAI_KEY = "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"

openai.api_key = OPENAI_KEY



def Gpt_answers(user_text):
    response = openai.ChatCompletion.create(
        model='gpt-3.5-turbo',
        messages=[
            {"role": "system", "content": "Hello, how can I assist you today?"}, 
            {"role": "user", "content": user_text}
        ],
        
    )

    return response.choices[0].message.content


# for i in range(1, 6):
# print(response)


# import openai

# # Подставьте свой ключ OpenAI API
# API_KEY = "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"

# # Текстовое описание
# PROMPT = "космический корабль летит над лазурной планетой"

# # Параметры генерации
# params = {
#     "prompt": PROMPT,
#     "n": 1,
#     "size": "512x512",
#     "response_format": "b64_json"
# }

# # Настройка OpenAI API
# openai.api_key = API_KEY

# # Генерация изображения
# response = openai.Image.create(**params)

# # Обработка ответа
# if response.status_code == 200:
#     data = json.loads(response.content)
    
#     # Получение изображения в формате Base64
#     image_b64 = data["data"][0]["b64_json"]
    
#     # Сохранение изображения
#     with open("image.png", "wb") as f:
#         f.write(base64.b64decode(image_b64))
    
#     # Вывод сообщения
#     print("Изображение успешно сгенерировано!")
# else:
#     print(f"Ошибка: {response.status_code}")


from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# Создаем новую презентацию
prs = Presentation()

# Добавляем первый слайд
def first_slide(tit_text, name_text):
    slide_layout = prs.slide_layouts[0]  # Выбираем макет слайда (5 - это макет с заголовком и содержанием)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"Презентация на тему\n{tit_text}"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Выравнивание текста по центру
    title.text_frame.paragraphs[0].font.size = Pt(24)  # Размер шрифта 24

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = f"Презентация создана\n{name_text}"

# Добавляем второй слайд
def second_slide(sc_t, th_t):
    slide_layout = prs.slide_layouts[5]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "План презентации"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на второй слайд
    text_contents = ["Вступление", sc_t, th_t, "Итог"]
    left = top = width = height = Inches(1)
    for numeratic, text_content in zip(range(1, 5), text_contents):
        txBox = slide.shapes.add_textbox(left, top+top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"{numeratic}. {text_content}"
        p.font.size = Pt(30)
        top += Inches(0.5)  # Устанавливаем расстояние между текстами

# Добавляем третий слайд
def third_slide(sub3_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Вступление"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub3_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def fourth_slide(sub4_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Вступление"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub4_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def fiveth_slide(sub5_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Вступление"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub5_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def sixth_slide(sub6_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Вступление"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub6_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)

def eighth_slide(sc_t, sub8_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub8_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)

def nineth_slide(sc_t, sub9_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub9_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def tenth_slide(sc_t, sub10_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub10_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def eleventh_slide(sc_t, sub11_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub11_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def twelveth_slide(sc_t, sub12_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub12_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def thirteenth_slide(sc_t, sub13_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub13_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def fourteenth_slide(sc_t, sub14_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub14_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def fiveteenth_slide(sc_t, sub15_text):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = sc_t
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  
    title.text_frame.paragraphs[0].font.size = Pt(30)  

    # Добавляем тексты на третий слайд
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = sub15_text
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(22)


def save_pr():
    first_slide("Франция в 19 веке", "Жавохир")
    second_slide("Деятельность", "История")
    user_input_text = input("Введите текст: ")
    response = Gpt_answers(user_input_text)
    print(response)
    sentence = response.split(".")
    third_slide(f"{sentence[0].strip()}. {sentence[1].strip()}. {sentence[2].strip()}. {sentence[3].strip()}.")
    fourth_slide(f"{sentence[4].strip()}. {sentence[5].strip()}. {sentence[6].strip()}. {sentence[7].strip()}.")
    fiveth_slide(f"{sentence[8].strip()}. {sentence[9].strip()}. {sentence[10].strip()}. {sentence[11].strip()}.")
    sixth_slide(f"{sentence[12].strip()}. {sentence[13].strip()}. {sentence[14].strip()}. {sentence[15].strip()}.")
    user_input_text2 = input("Введите текст: ")
    response2 = Gpt_answers(user_input_text2)
    print(response2)
    sentence2 = response2.split(".")
    eighth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    nineth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    tenth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    eleventh_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    twelveth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    thirteenth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    fourteenth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")
    fiveteenth_slide("Биография", "Let me introduce myself. My name is Mariya I am a 20-year-old student from Donetsk. I study at the university in my native town and my future profession is bookkeeping. I live with my parents and my elder sister Lena. We are a friendly family. Lena is 2 years older than me. We share our room and tell all our secrets to each other.")


save_pr()
prs.save("bot_for_b\presentation_test.pptx")